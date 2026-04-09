"""
혜화 데이터랩 디지털 사회와 나 — Lesson 1~6 PPT + PDF 생성기
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from fpdf import FPDF
import os

OUT = os.path.dirname(os.path.abspath(__file__))
FONT_REGULAR = "C:/Windows/Fonts/malgun.ttf"
FONT_BOLD    = "C:/Windows/Fonts/malgunbd.ttf"

# ─── 색상 팔레트 ──────────────────────────────────────────────
MOD_COLORS = {
    1: {"bg": RGBColor(0xF0, 0xFD, 0xF4), "accent": RGBColor(0x16, 0xA3, 0x4A),
        "title_bg": RGBColor(0x16, 0xA3, 0x4A), "hex": (0x16, 0xA3, 0x4A)},
    2: {"bg": RGBColor(0xEF, 0xF6, 0xFF), "accent": RGBColor(0x25, 0x63, 0xEB),
        "title_bg": RGBColor(0x25, 0x63, 0xEB), "hex": (0x25, 0x63, 0xEB)},
    3: {"bg": RGBColor(0xFA, 0xF5, 0xFF), "accent": RGBColor(0x93, 0x33, 0xEA),
        "title_bg": RGBColor(0x93, 0x33, 0xEA), "hex": (0x93, 0x33, 0xEA)},
}

# ─── 수업 내용 정의 ───────────────────────────────────────────
LESSONS = [
  {
    "id": "1-1", "module": 1,
    "title": "디지털 기술과 사회 변화",
    "subtitle": "Module 1 · 디지털 사회와 진로",
    "objectives": [
      "디지털 기술(AI, IoT, 빅데이터, 클라우드)의 개념을 설명할 수 있다",
      "사회가 농경→산업→디지털 사회로 변화해온 흐름을 이해한다",
      "디지털 격차의 개념과 실제 사례를 설명할 수 있다",
    ],
    "slides": [
      {
        "title": "디지털 기술이란?",
        "tag": "개념",
        "body": [
          "[TIP] 정보를 0과 1의 디지털 형식으로 처리하는 기술",
          "· 텍스트, 사진, 영상, 소리 모두 디지털 데이터로 변환",
          "· 핵심 키워드: 인터넷 · 스마트폰 · AI · IoT · 클라우드 · 빅데이터",
        ],
      },
      {
        "title": "사회 변화 타임라인",
        "tag": "개념",
        "body": [
          "[농경] 농경 사회 (~18세기)",
          "   땅 중심, 정보는 입에서 입으로, 변화 속도 느림",
          "",
          "[산업] 산업 사회 (18~20세기)",
          "   기계·대량 생산, 철도·전화·라디오로 정보 이동 빨라짐",
          "",
          "[디지털] 디지털 사회 (21세기~)",
          "   모든 것이 디지털로 연결, 데이터가 새로운 자원",
        ],
      },
      {
        "title": "지금 세상을 바꾸는 4가지 기술",
        "tag": "개념",
        "body": [
          "[AI] 인공지능(AI) — 컴퓨터가 스스로 학습하고 판단",
          "   예) 유튜브 추천, 챗봇, 얼굴 인식 잠금",
          "",
          "[IoT] 사물인터넷(IoT) — 일상 사물이 인터넷으로 연결",
          "   예) 스마트 냉장고, 스마트워치, 무인 결제기",
          "",
          "[데이터] 빅데이터 — 대량 데이터에서 패턴 발견",
          "   예) 넷플릭스 추천, 배달앱 인기 메뉴",
          "",
          "[클라우드] 클라우드 — 인터넷으로 저장공간·소프트웨어 제공",
          "   예) 구글 드라이브, 카카오톡 사진 백업",
        ],
      },
      {
        "title": "4가지 기술이 함께 작동하는 예",
        "tag": "개념",
        "body": [
          "[의료] 스마트워치(IoT)가 심박수 측정",
          "   → 빅데이터 분석 → AI가 이상 징후 판단 → 클라우드 기록",
          "",
          "[쇼핑] 편의점 카메라(IoT)로 고객 동선 수집",
          "   → 빅데이터 분석 → AI가 진열 추천 → 클라우드 저장",
          "",
          "[음악] 내 청취 기록(빅데이터)을 AI가 분석",
          "   → 클라우드에서 개인화 플레이리스트 제공",
        ],
      },
      {
        "title": "디지털 사회의 3가지 특징",
        "tag": "개념",
        "body": [
          "[연결] 온오프라인 경계 소멸",
          "   편의점 앱 주문, 온·오프라인 수업 혼합",
          "",
          "[메타] 사람과 사물의 네트워크",
          "   스마트 냉장고가 유통기한을 알려주고, 스마트워치가 심박수 측정",
          "",
          "[데이터] 빅데이터 생성",
          "   우리가 클릭·이동할 때마다 엄청난 데이터가 쌓임",
        ],
      },
      {
        "title": "디지털 격차 (Digital Divide)",
        "tag": "개념",
        "body": [
          "[주의] 디지털 기술을 사용할 수 있는 사람과 그렇지 못한 사람 사이의 정보 접근성 차이",
          "",
          "영향받는 그룹:",
          "[고령자] 고령자 — 새로운 기기·인터페이스에 익숙해지기 어려움",
          "[장애인] 장애인 — 기기 설계가 신체적 접근을 고려하지 않는 경우 多",
          "[저소득] 저소득층 — 최신 기기·인터넷 환경 갖추기 어려움",
          "",
          "[확인] 해결 방법: 음성 안내 / 큰 글씨 옵션 / 낮은 키오스크 / 직원 호출 버튼",
        ],
      },
      {
        "title": "활동: 키오스크 사용 경험 비교",
        "tag": "활동",
        "body": [
          "같은 키오스크인데 왜 사람마다 경험이 다를까요?",
          "",
          "[일반] 20대 일반 이용자",
          "   편리한 점: 빠른 주문, 메뉴 비교, 간편 결제",
          "   불편한 점: 오작동 시 당황",
          "",
          "[고령자] 70대 할머니",
          "   불편한 점: 작은 글씨, 시간 제한, 높은 화면, 직원 없음",
          "",
          "[장애인] 휠체어 이용자",
          "   불편한 점: 화면이 너무 높아 조작 불가, 낮은 키오스크 드묾",
          "",
          "[댓글] 생각해보기: 더 많은 사람이 편리하게 쓰려면 어떻게 개선해야 할까요?",
        ],
      },
      {
        "title": "정리",
        "tag": "정리",
        "body": [
          "✓ 디지털 기술 = 정보를 0과 1로 처리하는 기술",
          "✓ 사회 변화: 농경 → 산업 → 디지털 사회",
          "✓ 4대 기술: AI · IoT · 빅데이터 · 클라우드",
          "✓ 디지털 사회 특징: 온오프라인 융합, 사물 연결, 빅데이터 생성",
          "✓ 디지털 격차: 기술 접근성 차이 → 사회적 불평등 문제",
          "✓ 좋은 기술은 더 많은 사람이 쓸 수 있도록 설계되어야 해요",
        ],
      },
    ],
  },
  {
    "id": "1-2", "module": 1,
    "title": "디지털 사회의 직업과 진로",
    "subtitle": "Module 1 · 디지털 사회와 진로",
    "objectives": [
      "디지털 기술로 인해 사라지는/새로 생긴/변화하는 직업을 구분할 수 있다",
      "AI 트레이너, 데이터 분석가 등 새로운 디지털 직업을 설명할 수 있다",
      "자신의 디지털 역량을 자가진단하고 미래 진로와 연결할 수 있다",
    ],
    "slides": [
      {
        "title": "줄어드는 직업",
        "tag": "개념",
        "body": [
          "[계산원] 계산원 — 키오스크·무인 계산대·앱 결제 확산으로 수요 감소",
          "[잠금] 경비원 — CCTV·스마트 잠금장치·AI 영상 분석으로 자동화",
          "[전화] 텔레마케터 — AI 챗봇·자동 응답 시스템이 단순 상담 대체",
          "[여행] 여행사 오프라인 상담원 — 여행 플랫폼 앱·AI 추천 서비스 확산",
          "",
          "공통점: 단순 반복 업무 → 자동화에 취약",
        ],
      },
      {
        "title": "새로 생긴 직업",
        "tag": "개념",
        "body": [
          "[AI] AI 트레이너 — AI 모델 학습·교정, 편향 감시",
          "[데이터] 데이터 분석가 — 데이터로 패턴 발견·의사결정 지원",
          "[디자인] UX 디자이너 — 앱·웹의 사용자 경험 설계",
          "[보안] 사이버 보안 전문가 — 해킹·유출 방어 시스템 구축",
          "[메타] 메타버스 크리에이터 — 3D 가상 공간·아바타·디지털 아이템 제작",
          "[영상] 동영상 크리에이터 — 유튜브·틱톡 플랫폼 성장과 함께 등장",
        ],
      },
      {
        "title": "달라지는 직업들",
        "tag": "개념",
        "body": [
          "직업 자체는 남아 있지만 일하는 방식이 크게 변화",
          "",
          "[요리] 요리 연구가 + 생성형 AI → 조리법 초안을 AI가 도와줌",
          "[트레이너] 스포츠 트레이너 + IoT → 선수 몸에 센서 부착, 실시간 데이터 수집",
          "[디자인] 웹툰 작가 + 생성형 AI → AI가 채색·배경 작업 지원",
          "[농업] 과수 재배자 + IoT → 토양·날씨 센서로 스마트 농업",
          "[작가] 작가 + 생성형 AI → AI로 초안·아이디어 뽑고, 작가가 완성",
        ],
      },
      {
        "title": "활동: 직업 변화 분류해보기",
        "tag": "활동",
        "body": [
          "다음 직업들을 '줄어드는 직업' 또는 '새로 생긴 직업'으로 분류해보세요.",
          "",
          "① 계산원    ② 동영상 크리에이터",
          "③ 텔레마케터   ④ 빅데이터 분석가",
          "⑤ 경비원    ⑥ AI 트레이너",
          "⑦ 여행사 상담원  ⑧ 메타버스 크리에이터",
          "",
          "정답: 줄어드는 직업 → ①③⑤⑦  /  새로 생긴 직업 → ②④⑥⑧",
        ],
      },
      {
        "title": "활동: 기술과 직업 연결하기",
        "tag": "활동",
        "body": [
          "아래 기술과 직업을 올바르게 연결해보세요.",
          "",
          "기술                직업 + 활용",
          "빅데이터    →    스포츠 트레이너 (선수 데이터 수집·분석)",
          "인공지능(AI)  →    의사 (AI로 질병 정확하게 진단)",
          "클라우드    →    데이터 분석가 (대용량 데이터 저장·처리)",
          "생성형 AI   →    변호사 (법 조항 검색·서류 작성 자동화)",
        ],
      },
      {
        "title": "나의 디지털 역량 자가진단",
        "tag": "활동",
        "body": [
          "다음 항목 중 해당하는 것에 체크해보세요:",
          "",
          "□ 인터넷에서 정보를 빠르게 찾고 신뢰도를 판단할 수 있어요",
          "□ 프로그래밍·코딩의 기본 개념을 알고 있어요",
          "□ 표·그래프를 읽고 데이터에서 의미를 찾을 수 있어요",
          "□ 개인정보 보호, 피싱 등 기본 보안 수칙을 지키고 있어요",
          "□ 구글 독스·노션 같은 온라인 협업 도구를 써봤어요",
          "□ 사진·영상·글 등 디지털 콘텐츠를 직접 만들어봤어요",
          "□ ChatGPT, 제미나이 같은 AI 도구를 써봤어요",
          "",
          "5개 이상: 탄탄한 디지털 역량! / 3~4개: 기본 갖춤 / 2개 이하: 지금부터 시작!",
        ],
      },
      {
        "title": "정리",
        "tag": "정리",
        "body": [
          "✓ 줄어드는 직업: 단순 반복 업무 (계산원·경비원·텔레마케터 등)",
          "✓ 새로 생긴 직업: AI 트레이너·데이터 분석가·UX 디자이너·메타버스 크리에이터",
          "✓ 달라지는 직업: 기존 직업이 디지털 도구를 활용해 업그레이드",
          "✓ 핵심: 기술 자체보다 기술을 활용하는 능력이 중요",
          "✓ 미래 역량: 데이터 리터러시·AI 활용·디지털 보안 의식",
        ],
      },
    ],
  },
  {
    "id": "2-1", "module": 2,
    "title": "보호해야 할 정보와 공유해야 할 정보",
    "subtitle": "Module 2 · 정보 보호와 정보 공유",
    "objectives": [
      "정보 보호와 정보 공유의 개념과 차이를 설명할 수 있다",
      "개인정보의 정의와 유형을 이해하고 사례를 구분할 수 있다",
      "SNS에서 무심코 공개한 정보의 위험성을 인식할 수 있다",
    ],
    "slides": [
      {
        "title": "생각해보기: 수아의 이야기",
        "tag": "도입",
        "body": [
          "수아는 인스타그램 프로필에 학교 이름, 사진, 동네 이름을 올렸어요.",
          "틱톡에는 '서울 강북구, 고2, 매일 7시 등교'라고 소개했어요.",
          "",
          "어느 날 처음 보는 계정에서 DM이 왔어요.",
          "'OO고 2학년 수아 맞죠?'",
          "",
          "[댓글] 수아는 왜 이런 일이 생겼을까요?",
          "[댓글] 나라면 어떤 느낌이 들까요?",
        ],
      },
      {
        "title": "정보 보호 vs 정보 공유",
        "tag": "개념",
        "body": [
          "[보안] 정보 보호",
          "   위협으로부터 정보가 유출·변경·파괴되지 않도록 지키는 것",
          "   예) 군사 기밀 · 기업 기술 · 개인 정보",
          "",
          "[메타] 정보 공유",
          "   공공의 이익을 위해 정보를 알리거나 나누는 것",
          "   예) 날씨·재난 정보 · 진로·진학 정보 · 공공 데이터",
          "",
          "[주의] 한 번 공유된 정보는 되돌리기 어려워요!",
        ],
      },
      {
        "title": "개인정보란?",
        "tag": "개념",
        "body": [
          "살아있는 개인에 관한 정보로,",
          "그 정보만으로 또는 다른 정보와 결합했을 때",
          "특정 개인을 식별할 수 있는 것",
          "",
          "[신분증] 이름·주민번호    [전화] 전화번호·주소",
          "[의료] 진료 기록      [금융] 금융 정보",
          "[사진] 얼굴 사진      [위치] 위치 정보",
          "",
          "핵심: 단독으로 몰라도 조합하면 개인 식별 가능 → 개인정보!",
        ],
      },
      {
        "title": "활동: 이것도 개인정보일까?",
        "tag": "활동",
        "body": [
          "각 항목이 개인정보인지 판단해보세요:",
          "",
          "① '서울 거주, 고2, 운동 좋아함' (인스타 프로필)",
          "   → 개인정보 O  (사진·위치 결합 시 식별 가능)",
          "",
          "② '나 OO고등학교 다녀' (단톡방)",
          "   → 개인정보 O  (학년·반·사진과 결합 시 특정 가능)",
          "",
          "③ 스포티파이 공개 플레이리스트",
          "   → 개인정보 △  (스스로 공개한 취향, 단독으로는 비식별)",
          "",
          "④ SNS 게시물 위치 태그 '강남역 근처'",
          "   → 개인정보 O  (이동 경로 축적 → 생활 패턴 파악)",
        ],
      },
      {
        "title": "정보가 유출되면?",
        "tag": "개념",
        "body": [
          "전화번호 + 이름 + 학교 + 위치가 유출된다면...",
          "",
          "① 스팸·보이스피싱 폭발",
          "② 타깃 피싱: '선생님인 척', '학교 공지인 척' 속임",
          "③ 스토킹 위험: 집 근처 위치 + 귀가 시간 결합",
          "④ 명의 도용: 이름·주소로 택배 사기",
          "",
          "[주의] 한 번 유출된 정보는 인터넷 어딘가에 계속 남아요!",
        ],
      },
      {
        "title": "활동: 보호 vs 공유 분류",
        "tag": "활동",
        "body": [
          "아래 정보를 '보호해야 할 정보' / '공유해야 할 정보'로 나눠보세요:",
          "",
          "① 군사 기밀 정보        → 보호",
          "② 학교 급식 알레르기 정보  → 공유",
          "③ 처방전·진료 기록      → 보호",
          "④ 내일 서울 날씨 예보    → 공유",
          "⑤ 집 주소와 귀가 시간    → 보호",
          "⑥ 학교 공식 행사 일정    → 공유",
          "⑦ SNS 계정 아이디·비밀번호 → 보호",
          "⑧ 여행 맛집 블로그 후기   → 공유",
        ],
      },
      {
        "title": "초상권 이야기",
        "tag": "개념",
        "body": [
          "유튜버 A가 식당에서 촬영하면서 주변 손님 얼굴이 화면에 담겼어요.",
          "'공개 장소니까 괜찮다'고 생각했어요.",
          "",
          "❓ A의 생각이 맞을까요?",
          "",
          "답: 아니요! — 초상권 침해에 해당할 수 있어요.",
          "",
          "초상권 = 자신의 얼굴·모습이 허락 없이 촬영·공개되지 않을 권리",
          "공개 장소라도 얼굴이 찍혀 동의 없이 공개되면 침해!",
          "→ 촬영 전 동의 받기, 또는 편집 시 얼굴 가리기",
        ],
      },
      {
        "title": "정리",
        "tag": "정리",
        "body": [
          "✓ 정보 보호: 위협으로부터 정보를 지키는 것",
          "✓ 정보 공유: 공공의 이익을 위해 알려야 하는 정보",
          "✓ 개인정보: 단독 또는 조합으로 개인을 식별할 수 있는 정보",
          "✓ 한 번 유출된 정보는 되돌리기 어려워요",
          "✓ SNS에 올리는 위치·학교·패턴 조합에 주의해요",
          "✓ 초상권: 공개 장소에서도 사전 동의 없는 촬영·공개는 불법",
        ],
      },
    ],
  },
  {
    "id": "2-2", "module": 2,
    "title": "올바른 정보 보호 방법",
    "subtitle": "Module 2 · 정보 보호와 정보 공유",
    "objectives": [
      "스팸·피싱·해킹·랜섬웨어 등 디지털 위협의 종류를 구분할 수 있다",
      "안전한 비밀번호 설정 기준을 적용할 수 있다",
      "2단계 인증의 개념과 필요성을 설명할 수 있다",
    ],
    "slides": [
      {
        "title": "디지털 위협의 종류",
        "tag": "개념",
        "body": [
          "[이메일] 스팸 (Spam)",
          "   불특정 다수에게 보내는 광고성·불필요한 메시지",
          "   대처: 수신 거부, 모르는 링크 클릭 금지",
          "",
          "[피싱] 피싱 (Phishing)",
          "   신뢰 기관인 척 속여 개인정보·금융정보 빼내기",
          "   예) '계좌 정지됩니다' + 가짜 사이트 링크",
          "",
          "[디지털] 해킹 (Hacking)",
          "   허락 없이 타인 시스템 침입해 정보 열람·유출",
          "",
          "[잠금] 랜섬웨어 (Ransomware)",
          "   파일을 암호화해 사용 불가 → 돈 요구",
        ],
      },
      {
        "title": "안전한 비밀번호 만들기",
        "tag": "개념",
        "body": [
          "[확인] 좋은 비밀번호의 조건:",
          "   · 8자리 이상 (12자 이상 권장)",
          "   · 영문 대소문자 + 숫자 + 특수문자 조합",
          "   · 생일·이름·학번 등 개인정보 포함 금지",
          "   · 여러 사이트에 같은 비밀번호 사용 금지",
          "   · 3~6개월마다 주기적으로 변경",
          "",
          "❌ 위험한 비밀번호 예시:",
          "   12345678 / 000000 / abc1234 / 생년월일",
        ],
      },
      {
        "title": "활동: 비밀번호 강도 판단",
        "tag": "활동",
        "body": [
          "아래 비밀번호를 '강함' / '약함'으로 분류해보세요:",
          "",
          "① 19900101          → 약함 (생년월일, 단순 숫자)",
          "② Hana#2024!Sky     → 강함 (대소문자+숫자+특수문자+12자 이상)",
          "③ password123       → 약함 (흔한 패턴)",
          "④ !Blue7Tree*Lamp   → 강함 (무작위 단어+특수문자 조합)",
          "⑤ 000000            → 약함 (단순 반복)",
          "⑥ k9@mL#2pQr!       → 강함 (복잡한 조합)",
        ],
      },
      {
        "title": "2단계 인증이란?",
        "tag": "개념",
        "body": [
          "비밀번호(1단계) 입력 후,",
          "추가로 본인 확인을 하는 방법",
          "",
          "예시:",
          "[스마트폰] 스마트폰으로 문자 인증 코드 받기",
          "[이메일] 이메일로 확인 코드 받기",
          "[열쇠] 인증 앱(구글 OTP 등)으로 코드 생성",
          "",
          "왜 필요할까?",
          "→ 비밀번호가 유출돼도 추가 방어막이 생겨요",
          "→ 해커가 비밀번호를 알아도 내 스마트폰이 없으면 로그인 불가",
        ],
      },
      {
        "title": "피싱 문자 구별하기",
        "tag": "활동",
        "body": [
          "다음 중 피싱 문자는 무엇일까요?",
          "",
          "① '[국세청] 환급금 128,000원이 미수령입니다. 아래 링크에서 신청하세요: bit.ly/xxx'",
          "   → 피싱! (공식 기관은 bit.ly 단축링크 사용 안 함)",
          "",
          "② '[카카오뱅크] 비정상 로그인 감지. 즉시 확인 필요: kakao-secure.com'",
          "   → 피싱! (공식 도메인은 kakaobank.com)",
          "",
          "③ '[건강보험공단] 건강검진 대상자 안내. 자세한 내용은 nhis.or.kr에서 확인'",
          "   → 안전 (공식 도메인 nhis.or.kr 안내)",
        ],
      },
      {
        "title": "나의 보안 습관 점검",
        "tag": "활동",
        "body": [
          "□ 비밀번호가 8자리 이상이에요",
          "□ 여러 사이트에 다른 비밀번호를 사용해요",
          "□ 중요 계정에 2단계 인증을 설정했어요",
          "□ 모르는 링크·첨부파일은 클릭하지 않아요",
          "□ 공공 와이파이에서 금융 거래를 하지 않아요",
          "□ 중요한 파일을 정기적으로 백업해요",
          "",
          "5~6개: 보안 습관 우수! / 3~4개: 부분 개선 필요 / 2개 이하: 지금 바로 점검!",
        ],
      },
      {
        "title": "정리",
        "tag": "정리",
        "body": [
          "✓ 스팸: 대량 광고성 메시지 → 수신 거부, 링크 클릭 금지",
          "✓ 피싱: 기관 사칭 → 공식 앱·홈페이지에서 직접 확인",
          "✓ 해킹: 시스템 침입 → 강력한 비밀번호 + 보안 업데이트",
          "✓ 랜섬웨어: 파일 암호화 → 정기 백업, 출처 불명 파일 실행 금지",
          "✓ 좋은 비밀번호: 12자 이상 + 대소문자+숫자+특수문자 조합",
          "✓ 2단계 인증: 비밀번호 유출돼도 추가 방어막",
        ],
      },
    ],
  },
  {
    "id": "3-1", "module": 3,
    "title": "정보 보안의 개념과 필요성",
    "subtitle": "Module 3 · 정보 보안",
    "objectives": [
      "정보 보안의 3대 요소(기밀성·무결성·가용성)를 설명할 수 있다",
      "4가지 위협 요소(훼손·변조·접근유출·위조)를 구분할 수 있다",
      "피싱 문자를 보고 실제 위협 상황을 판단할 수 있다",
    ],
    "slides": [
      {
        "title": "정보 보안이란?",
        "tag": "개념",
        "body": [
          "정보 보안 (Information Security)",
          "",
          "정보를 허가되지 않은 접근·변경·공개·파괴로부터",
          "보호하는 과정과 방법론",
          "",
          "왜 필요할까?",
          "[스마트폰] 스마트폰에는 연락처·사진·금융 앱이 모두 연결",
          "[의료] 병원 의료 시스템이 해킹당하면 환자 생명까지 위험",
          "[정부] 정부 데이터가 유출되면 국가 안보 위협",
          "",
          "→ 디지털 기술이 발전할수록 보안의 중요성도 커진다!",
        ],
      },
      {
        "title": "4가지 위협 요소",
        "tag": "개념",
        "body": [
          "[훼손] 훼손 — 정보를 손상·파괴해 정상 작동 불가",
          "   예) 바이러스가 중요 파일을 손상시켜 열 수 없게 만듦",
          "",
          "[변조] 변조 — 정보 내용을 다른 내용으로 바꿈",
          "   예) 성적 데이터 몰래 수정, 뉴스 기사 변경",
          "",
          "[접근] 접근과 유출 — 허락 없이 정보 확인·외부 유출",
          "   예) 퇴사 직원이 회사 고객 정보 외부에 넘김",
          "",
          "[위조] 위조 — 거짓 정보를 정상 정보인 척 속임",
          "   예) 학교에서 오는 척 꾸민 가짜 이메일",
        ],
      },
      {
        "title": "정보 보안 3대 요소 (CIA)",
        "tag": "개념",
        "body": [
          "[보안] 기밀성 (Confidentiality)",
          "   허락된 사용자만 정보에 접근 가능",
          "   예) 성적 정보는 해당 학생과 담임 선생님만 열람",
          "",
          "[확인] 무결성 (Integrity)",
          "   정보를 허락 없이 변경할 수 없음",
          "   예) 성적표 점수가 누군가에 의해 바뀌지 않아야 함",
          "",
          "⚡ 가용성 (Availability)",
          "   사용자가 원할 때 언제든지 정보에 접근 가능",
          "   예) 온라인 시험 중 서버 다운 = 가용성 침해",
        ],
      },
      {
        "title": "활동: 보안 위협 분류",
        "tag": "활동",
        "body": [
          "다음을 '스팸·피싱 위협' / '해킹·랜섬웨어 위협'으로 분류하세요:",
          "",
          "① 수신 거부 설정, 모르는 링크 클릭 금지",
          "   → 스팸·피싱 대처",
          "② 중요 파일 외장하드·클라우드 백업",
          "   → 해킹·랜섬웨어 대처",
          "③ 공식 앱·홈페이지에서 직접 확인",
          "   → 스팸·피싱 대처",
          "④ 출처 불명 파일 실행 금지",
          "   → 해킹·랜섬웨어 대처",
        ],
      },
      {
        "title": "활동: 3대 요소 매칭",
        "tag": "활동",
        "body": [
          "다음 상황이 어떤 보안 요소를 침해하는지 연결해보세요:",
          "",
          "① 해커가 내 카카오톡 메시지를 읽음",
          "   → 기밀성 침해",
          "",
          "② 랜섬웨어가 파일을 암호화해 열 수 없게 됨",
          "   → 가용성 침해",
          "",
          "③ 악성코드가 학교 성적 데이터를 몰래 수정",
          "   → 무결성 침해",
        ],
      },
      {
        "title": "실천 방법",
        "tag": "개념",
        "body": [
          "[열쇠] 강력한 비밀번호 + 2단계 인증",
          "   개인 정보·생일 포함 금지, 주기적 변경",
          "",
          "[업데이트] 소프트웨어 업데이트",
          "   보안 취약점 패치 → 해킹 경로 차단",
          "",
          "[백업] 정기 백업",
          "   외장하드·클라우드에 중요 파일 백업",
          "   랜섬웨어 피해 최소화",
        ],
      },
      {
        "title": "정리",
        "tag": "정리",
        "body": [
          "✓ 정보 보안 = 정보를 허가되지 않은 접근·변경·파괴로부터 보호",
          "✓ 위협 4가지: 훼손 · 변조 · 접근·유출 · 위조",
          "✓ 3대 요소: 기밀성(접근제한) · 무결성(변경금지) · 가용성(언제든 접근)",
          "✓ 실천: 강력한 비밀번호 + 업데이트 + 백업",
          "✓ 의심스러운 링크·파일은 절대 클릭·실행 금지",
        ],
      },
    ],
  },
  {
    "id": "3-2", "module": 3,
    "title": "디지털 윤리와 안전한 디지털 사회",
    "subtitle": "Module 3 · 정보 보안",
    "objectives": [
      "디지털 윤리의 개념과 주요 문제를 설명할 수 있다",
      "사이버 폭력의 유형을 구분하고 올바른 대처 방법을 말할 수 있다",
      "일상 속 디지털 윤리 딜레마 상황을 판단할 수 있다",
    ],
    "slides": [
      {
        "title": "디지털 윤리란?",
        "tag": "개념",
        "body": [
          "디지털 공간에서 지켜야 할 도덕적 기준과 책임",
          "",
          "주요 문제들:",
          "[댓글] 악성 댓글 — 익명성 뒤에 숨어 상처를 주는 댓글",
          "[복사] 개인정보·저작권 침해 — 허락 없이 사진·글·음악 무단 사용",
          "[뉴스] 가짜 뉴스 전파 — 사실 확인 없이 자극적 정보 공유",
          "[스마트폰] 디지털 중독 — 과도한 스마트폰·게임·SNS 사용",
          "",
          "온라인이라도 현실과 똑같은 사람 사이의 관계예요!",
        ],
      },
      {
        "title": "사이버 폭력의 종류",
        "tag": "개념",
        "body": [
          "[분노] 사이버 언어폭력",
          "   욕설·모욕·협박을 온라인으로 보내는 것",
          "   예) 게임에서 진다고 상대방에게 심한 욕설 DM",
          "",
          "[따돌림] 사이버 따돌림 (카따)",
          "   단톡방에서 특정인을 의도적으로 무시·배제",
          "   예) 반 단톡방에서 한 친구만 읽씹하기로 몰래 약속",
          "",
          "[사진] 사이버 스토킹",
          "   온라인에서 타인을 집요하게 추적·괴롭힘",
          "",
          "[이미지] 사이버 명예훼손",
          "   사실이든 아니든 명예 훼손 정보를 퍼뜨림",
        ],
      },
      {
        "title": "활동: 이런 행동, 괜찮을까?",
        "tag": "활동",
        "body": [
          "다음 상황에 대해 윤리적으로 생각해봐요:",
          "",
          "① 친구가 파티에서 찍힌 사진을 내 SNS에 올렸어요.",
          "   → 문제 있음! 친구의 동의 없이 얼굴 사진 공유는 초상권 침해",
          "",
          "② 친구와의 카카오톡 대화를 캡처해 다른 친구에게 보냈어요.",
          "   → 문제 있음! 사적 대화 내용은 개인정보, 무단 공유는 침해",
          "",
          "③ 예쁜 폰트를 인터넷에서 무료로 받아 발표 PPT에 사용했어요.",
          "   → 반드시 확인 필요! 유료 폰트 무단 사용은 저작권 침해",
        ],
      },
      {
        "title": "저작권 바로 알기",
        "tag": "개념",
        "body": [
          "저작권 = 창작물을 만든 사람이 갖는 권리",
          "",
          "저작권으로 보호되는 것:",
          "[글] 글 · [사진] 사진 · [음악] 음악 · [영상] 영상 · [디자인] 그림 · [디지털] 프로그램",
          "",
          "주의해야 할 상황:",
          "· 유튜브 영상에 저작권 있는 음악 사용",
          "· 인터넷에서 이미지 무단으로 가져다 발표 자료에 사용",
          "· 책 내용을 그대로 복사해 과제 제출",
          "",
          "올바른 방법:",
          "출처 표기 / 저작자 동의 / 무료 라이센스 자료 활용",
        ],
      },
      {
        "title": "활동: 윤리 vs 위반 분류",
        "tag": "활동",
        "body": [
          "아래 행동을 '올바른 행동' / '윤리 위반'으로 분류해보세요:",
          "",
          "① 친구 동의를 받고 사진을 SNS에 공유  → 올바른 행동",
          "② 무료 CC 라이선스 음악을 유튜브 영상에 사용 → 올바른 행동",
          "③ 게임에서 졌다고 상대방에게 욕설 메시지 → 윤리 위반",
          "④ 친구 몰래 SNS 대화 캡처해서 유포  → 윤리 위반",
          "⑤ 가짜 뉴스를 팩트체크 없이 단톡방 공유 → 윤리 위반",
          "⑥ 유료 폰트 무단 사용       → 윤리 위반",
          "⑦ 사이버 폭력 피해 친구를 위해 신고·위로 → 올바른 행동",
        ],
      },
      {
        "title": "나의 디지털 생활 돌아보기",
        "tag": "활동",
        "body": [
          "□ 다른 사람 동의 없이 사진·동영상을 공유하지 않아요",
          "□ 온라인에서도 상대방을 존중하는 말을 사용해요",
          "□ 음악·이미지·영상 사용 시 저작권을 확인해요",
          "□ 사실 확인 없이 정보를 공유하지 않아요",
          "□ 사이버 폭력 피해자를 방관하지 않아요",
          "□ 스마트폰 사용 시간을 스스로 조절해요",
          "□ 개인정보를 SNS에 무분별하게 공개하지 않아요",
          "",
          "6~7개: 디지털 시민 의식 우수! / 4~5개: 좋은 습관 중 / 3개 이하: 오늘부터 실천!",
        ],
      },
      {
        "title": "정리",
        "tag": "정리",
        "body": [
          "✓ 디지털 윤리: 디지털 공간에서 지켜야 할 도덕적 기준",
          "✓ 사이버 폭력 유형: 언어폭력·따돌림·스토킹·명예훼손",
          "✓ 초상권: 동의 없는 사진·영상 공개 금지",
          "✓ 저작권: 창작물 사용 시 반드시 출처 표기·동의",
          "✓ 가짜 뉴스: 공유 전 반드시 팩트체크",
          "✓ 디지털 시민 = 권리와 책임을 함께 갖는 사람",
        ],
      },
    ],
  },
]

# ─── PPT 생성 ─────────────────────────────────────────────────

def add_text_box(slide, text, left, top, width, height,
                 font_size=14, bold=False, color=None,
                 align=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    if color:
        run.font.color.rgb = color
    return txBox

def set_slide_bg(slide, r, g, b):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)

def make_cover_slide(prs, lesson, colors):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide, *colors["hex"])

    # 모듈 레이블
    add_text_box(slide, lesson["subtitle"], 0.4, 0.3, 9, 0.5,
                 font_size=14, color=RGBColor(0xFF, 0xFF, 0xFF),
                 font_name="맑은 고딕")

    # 레슨 번호
    add_text_box(slide, f"Lesson {lesson['id']}", 0.4, 0.8, 9, 0.7,
                 font_size=28, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                 font_name="맑은 고딕")

    # 레슨 타이틀
    add_text_box(slide, lesson["title"], 0.4, 1.5, 9, 1.2,
                 font_size=36, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                 font_name="맑은 고딕")

    # 하단 구분선 느낌
    add_text_box(slide, "혜화 데이터랩 · 디지털 사회와 나", 0.4, 6.8, 9, 0.4,
                 font_size=11, color=RGBColor(0xFF, 0xFF, 0xFF),
                 font_name="맑은 고딕")
    return slide

def make_objectives_slide(prs, lesson, colors):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide, 0xFF, 0xFF, 0xFF)

    # 상단 컬러 바
    from pptx.util import Inches, Emu
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(1.0)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = colors["accent"]
    bar.line.fill.background()

    add_text_box(slide, "학습 목표", 0.4, 0.15, 9, 0.6,
                 font_size=20, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                 font_name="맑은 고딕")

    body = "\n".join(f"✓  {obj}" for obj in lesson["objectives"])
    add_text_box(slide, body, 0.5, 1.2, 9, 4,
                 font_size=16, color=RGBColor(0x1F, 0x29, 0x37),
                 font_name="맑은 고딕")
    return slide

def make_content_slide(prs, slide_data, colors):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide, 0xFF, 0xFF, 0xFF)

    # 상단 바
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(1.0)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = colors["accent"]
    bar.line.fill.background()

    # 태그 + 타이틀
    tag = slide_data.get("tag", "")
    title = slide_data["title"]
    add_text_box(slide, f"[{tag}]  {title}", 0.4, 0.15, 9.2, 0.65,
                 font_size=18, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                 font_name="맑은 고딕")

    # 본문
    body_lines = slide_data.get("body", [])
    body_text = "\n".join(body_lines)
    add_text_box(slide, body_text, 0.5, 1.1, 9, 5.6,
                 font_size=14, color=RGBColor(0x1F, 0x29, 0x37),
                 font_name="맑은 고딕")
    return slide

def create_ppt(lessons, out_path):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    for lesson in lessons:
        colors = MOD_COLORS[lesson["module"]]
        make_cover_slide(prs, lesson, colors)
        make_objectives_slide(prs, lesson, colors)
        for slide_data in lesson["slides"]:
            make_content_slide(prs, slide_data, colors)

    prs.save(out_path)
    print(f"PPT 저장 완료: {out_path}")


# ─── PDF 생성 ─────────────────────────────────────────────────

class KoreanPDF(FPDF):
    def __init__(self):
        super().__init__()
        self.add_font("Malgun", "", FONT_REGULAR, uni=True)
        self.add_font("Malgun", "B", FONT_BOLD, uni=True)

    def header(self):
        pass

    def footer(self):
        self.set_y(-12)
        self.set_font("Malgun", "", 8)
        self.set_text_color(150, 150, 150)
        self.cell(0, 5, f"혜화 데이터랩 · 디지털 사회와 나  |  {self.page_no()}", align="C")

def pdf_cover(pdf, lesson, mod):
    c = MOD_COLORS[mod]
    r, g, b = c["hex"]
    pdf.add_page()
    pdf.set_fill_color(r, g, b)
    pdf.rect(0, 0, 210, 297, "F")

    # 서브타이틀
    pdf.set_y(60)
    pdf.set_font("Malgun", "", 12)
    pdf.set_text_color(255, 255, 255)
    pdf.cell(0, 8, lesson["subtitle"], align="C", ln=True)

    # 레슨 번호
    pdf.ln(8)
    pdf.set_font("Malgun", "B", 22)
    pdf.cell(0, 12, f"Lesson {lesson['id']}", align="C", ln=True)

    # 타이틀
    pdf.ln(4)
    pdf.set_font("Malgun", "B", 30)
    pdf.cell(0, 16, lesson["title"], align="C", ln=True)

    # 하단
    pdf.set_y(260)
    pdf.set_font("Malgun", "", 10)
    pdf.cell(0, 6, "혜화 데이터랩 · 디지털 사회와 나", align="C", ln=True)

def pdf_objectives(pdf, lesson, mod):
    c = MOD_COLORS[mod]
    r, g, b = c["hex"]
    pdf.add_page()

    # 상단 바
    pdf.set_fill_color(r, g, b)
    pdf.rect(0, 0, 210, 18, "F")
    pdf.set_y(4)
    pdf.set_font("Malgun", "B", 14)
    pdf.set_text_color(255, 255, 255)
    pdf.cell(0, 8, "  학습 목표", ln=True)

    pdf.ln(8)
    pdf.set_text_color(30, 30, 30)
    for obj in lesson["objectives"]:
        pdf.set_font("Malgun", "B", 12)
        pdf.set_x(15)
        pdf.cell(6, 8, "✓", ln=False)
        pdf.set_font("Malgun", "", 12)
        pdf.multi_cell(170, 8, obj)
        pdf.ln(2)

def pdf_content(pdf, slide_data, mod):
    c = MOD_COLORS[mod]
    r, g, b = c["hex"]
    pdf.add_page()

    # 상단 바
    pdf.set_fill_color(r, g, b)
    pdf.rect(0, 0, 210, 18, "F")
    pdf.set_y(3)
    tag = slide_data.get("tag", "")
    title = slide_data["title"]
    pdf.set_font("Malgun", "B", 13)
    pdf.set_text_color(255, 255, 255)
    pdf.cell(0, 10, f"  [{tag}]  {title}", ln=True)

    pdf.ln(5)
    pdf.set_text_color(30, 30, 30)
    for line in slide_data.get("body", []):
        if line == "":
            pdf.ln(3)
        else:
            pdf.set_x(15)
            pdf.set_font("Malgun", "", 11)
            pdf.multi_cell(180, 7, line)

def create_pdf(lessons, out_path):
    pdf = KoreanPDF()
    pdf.set_auto_page_break(auto=True, margin=15)

    for lesson in lessons:
        mod = lesson["module"]
        pdf_cover(pdf, lesson, mod)
        pdf_objectives(pdf, lesson, mod)
        for slide_data in lesson["slides"]:
            pdf_content(pdf, slide_data, mod)

    pdf.output(out_path)
    print(f"PDF 저장 완료: {out_path}")


# ─── 실행 ──────────────────────────────────────────────────────
if __name__ == "__main__":
    ppt_path = os.path.join(OUT, "디지털사회와나_Lesson1-6.pptx")
    pdf_path = os.path.join(OUT, "디지털사회와나_Lesson1-6.pdf")

    create_ppt(LESSONS, ppt_path)
    create_pdf(LESSONS, pdf_path)
    print("\n완료! materials 폴더를 확인하세요.")
